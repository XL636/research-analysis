"""PPTX 解析器测试."""
from __future__ import annotations

import pytest
from pptx import Presentation
from src.core.models import FileType
from src.parsers.ppt_parser import parse_pptx


@pytest.fixture
def sample_pptx(tmp_path):
    """Generate a sample PPTX file."""
    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()

    # Slide 1 - Title slide
    slide_layout = prs.slide_layouts[0]  # Title layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Machine Learning Overview"
    slide.placeholders[1].text = "A Comprehensive Study"

    # Slide 2 - Content slide
    slide_layout = prs.slide_layouts[1]  # Title + content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = "Machine learning is transforming many fields."

    # Slide 3 - Another content slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Results"
    slide.placeholders[1].text = "Our experiments show improved accuracy."

    prs.save(str(pptx_path))
    return str(pptx_path)


def test_parse_pptx_returns_parsed_document(sample_pptx):
    result = parse_pptx(sample_pptx)
    assert result.file_type == FileType.PPTX
    assert result.file_path == sample_pptx


def test_parse_pptx_extracts_title(sample_pptx):
    result = parse_pptx(sample_pptx)
    assert result.title == "Machine Learning Overview"


def test_parse_pptx_extracts_sections(sample_pptx):
    result = parse_pptx(sample_pptx)
    # Should have sections corresponding to slides
    assert len(result.sections) >= 2


def test_parse_pptx_extracts_raw_text(sample_pptx):
    result = parse_pptx(sample_pptx)
    assert "machine learning" in result.raw_text.lower()


def test_parse_pptx_file_not_found():
    with pytest.raises(FileNotFoundError):
        parse_pptx("/nonexistent/file.pptx")
