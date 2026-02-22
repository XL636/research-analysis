"""PPT 解析器 - 使用 python-pptx 提取幻灯片内容."""

from __future__ import annotations

from pathlib import Path

from loguru import logger
from pptx import Presentation

from src.core.models import DocumentSection, FileType, ParsedDocument


def parse_pptx(file_path: str) -> ParsedDocument:
    """解析 PPTX 文件，提取幻灯片文本和结构."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {file_path}")

    logger.info(f"Parsing PPTX: {path.name}")

    prs = Presentation(file_path)
    sections: list[DocumentSection] = []
    all_text_parts: list[str] = []
    title = ""

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        slide_title = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue
                slide_texts.append(text)

            # 尝试从标题占位符获取标题
            try:
                ph_format = shape.placeholder_format
                if ph_format is not None and ph_format.idx == 0:
                    slide_title = shape.text_frame.text.strip()
            except ValueError:
                pass  # 非占位符 shape，跳过

        # 第一页标题作为文档标题
        if slide_num == 1 and slide_title and not title:
            title = slide_title

        if slide_texts:
            content = "\n".join(slide_texts)
            section_title = slide_title or f"Slide {slide_num}"

            sections.append(DocumentSection(
                title=section_title,
                content=content,
                level=2,
            ))
            all_text_parts.append(f"[Slide {slide_num}] {section_title}\n{content}")

    raw_text = "\n\n".join(all_text_parts)

    # 提取元数据
    metadata: dict[str, str] = {}
    if prs.core_properties:
        props = prs.core_properties
        if props.title:
            title = title or props.title
            metadata["pptx_title"] = props.title
        if props.author:
            metadata["pptx_author"] = props.author

    logger.info(f"PPTX parsed: {len(sections)} slides, {len(raw_text)} chars")

    authors = [metadata["pptx_author"]] if "pptx_author" in metadata else []

    return ParsedDocument(
        file_path=file_path,
        file_type=FileType.PPTX,
        title=title,
        authors=authors,
        sections=sections,
        raw_text=raw_text,
        metadata=metadata,
    )
