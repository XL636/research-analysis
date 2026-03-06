"""阅读助手 - 文件分页提取服务."""

from __future__ import annotations

import re
import shutil
import uuid
from pathlib import Path

import pymupdf
from loguru import logger


UPLOAD_DIR = Path("./uploads/reader")
ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".md", ".txt"}


def save_reader_file(file_name: str, file_bytes: bytes) -> tuple[str, str]:
    """Save uploaded file, return (doc_id, file_path)."""
    doc_id = uuid.uuid4().hex[:12]
    dest_dir = UPLOAD_DIR / doc_id
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest = dest_dir / Path(file_name).name
    dest.write_bytes(file_bytes)
    return doc_id, str(dest)


def delete_reader_files(file_path: str) -> None:
    """Delete stored file directory."""
    p = Path(file_path)
    if p.exists():
        parent = p.parent
        if parent.name and parent != UPLOAD_DIR:
            shutil.rmtree(parent, ignore_errors=True)


def extract_pages(file_path: str, file_type: str) -> list[str]:
    """Extract text pages from file. Returns list of page texts."""
    ft = file_type.lower()
    if ft == "pdf":
        return _extract_pdf_pages(file_path)
    elif ft == "pptx":
        return _extract_pptx_pages(file_path)
    elif ft == "docx":
        return _extract_docx_pages(file_path)
    elif ft in ("md", "markdown"):
        return _extract_md_pages(file_path)
    else:  # txt and others
        return _extract_txt_pages(file_path)


def detect_file_type(file_name: str) -> str:
    """Detect file type from extension."""
    ext = Path(file_name).suffix.lower()
    mapping = {
        ".pdf": "pdf",
        ".pptx": "pptx",
        ".docx": "docx",
        ".md": "md",
        ".markdown": "md",
        ".txt": "txt",
    }
    return mapping.get(ext, "txt")


def _extract_pdf_pages(file_path: str) -> list[str]:
    """PDF: each natural page = one page."""
    doc = pymupdf.open(file_path)
    pages = []
    for page in doc:
        text = page.get_text("text")
        pages.append(text)
    doc.close()
    logger.debug(f"PDF extracted: {len(pages)} pages")
    return pages


def _extract_pptx_pages(file_path: str) -> list[str]:
    """PPTX: each slide = one page."""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        slide_title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
            try:
                ph = shape.placeholder_format
                if ph is not None and ph.idx == 0:
                    slide_title = shape.text_frame.text.strip()
            except ValueError:
                pass
        header = f"## Slide {slide_num}"
        if slide_title:
            header += f": {slide_title}"
        page_text = header + "\n\n" + "\n\n".join(texts) if texts else header
        pages.append(page_text)
    logger.debug(f"PPTX extracted: {len(pages)} slides")
    return pages


def _extract_docx_pages(file_path: str) -> list[str]:
    """DOCX: split by Headings; if no headings, chunk by ~2000 chars."""
    try:
        import docx
    except ImportError:
        logger.warning("python-docx not installed, falling back to txt")
        return _extract_txt_pages(file_path)

    doc = docx.Document(file_path)
    sections: list[str] = []
    current_lines: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = para.style.name if para.style else ""
        is_heading = "Heading" in style_name or "heading" in style_name

        if is_heading and current_lines:
            sections.append("\n".join(current_lines))
            current_lines = []

        if is_heading:
            current_lines.append(f"## {text}")
        else:
            current_lines.append(text)

    if current_lines:
        sections.append("\n".join(current_lines))

    # If no headings found (single section), chunk by ~2000 chars
    if len(sections) <= 1 and sections:
        full_text = sections[0]
        if len(full_text) > 2500:
            sections = _chunk_text(full_text, 2000)

    logger.debug(f"DOCX extracted: {len(sections)} pages")
    return sections


def _extract_md_pages(file_path: str) -> list[str]:
    """Markdown: split by ## headings; if none, chunk by paragraphs."""
    text = Path(file_path).read_text(encoding="utf-8")

    # Split by ## headings
    parts = re.split(r"(?=^## )", text, flags=re.MULTILINE)
    pages = [p.strip() for p in parts if p.strip()]

    if len(pages) <= 1 and pages:
        full_text = pages[0]
        if len(full_text) > 2500:
            pages = _chunk_text(full_text, 2000)

    logger.debug(f"MD extracted: {len(pages)} pages")
    return pages


def _extract_txt_pages(file_path: str) -> list[str]:
    """TXT: chunk by paragraphs merged to ~2000 chars."""
    text = Path(file_path).read_text(encoding="utf-8")
    if not text.strip():
        return [""]

    pages = _chunk_text(text, 2000)
    logger.debug(f"TXT extracted: {len(pages)} pages")
    return pages


def sample_document_content(store, doc_id: int, max_chars: int = 16000) -> str:
    """Sample document content for overview/study generation.

    Takes pages from beginning, middle, and end to get representative content.
    For longer documents (>20 pages), also samples at 1/4 and 3/4 positions.
    """
    from src.store.reader_store import ReaderStore
    doc = store.get_document(doc_id)
    if not doc:
        return ""
    total = doc["total_pages"]
    if total == 0:
        return ""

    # Strategy: first 3 pages, middle 2, last 2 (up to total)
    page_nums = set()
    # First pages
    for i in range(1, min(4, total + 1)):
        page_nums.add(i)
    # Middle pages
    mid = total // 2
    for i in range(max(1, mid - 1), min(total + 1, mid + 2)):
        page_nums.add(i)
    # Last pages
    for i in range(max(1, total - 1), total + 1):
        page_nums.add(i)
    # For longer documents, add 1/4 and 3/4 sampling points
    if total > 20:
        q1 = total // 4
        q3 = total * 3 // 4
        page_nums.add(max(1, q1))
        page_nums.add(min(total, q3))

    pages = store.get_pages_by_nums(doc_id, sorted(page_nums))
    parts = []
    char_count = 0
    for p in pages:
        text = p["content"]
        if char_count + len(text) > max_chars:
            text = text[:max_chars - char_count]
        parts.append(f"[第 {p['page_num']} 页]\n{text}")
        char_count += len(text)
        if char_count >= max_chars:
            break

    return "\n\n---\n\n".join(parts)


def _chunk_text(text: str, chunk_size: int = 2000) -> list[str]:
    """Split text into chunks of approximately chunk_size characters at paragraph boundaries."""
    paragraphs = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    current: list[str] = []
    current_len = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if current_len + len(para) > chunk_size and current:
            chunks.append("\n\n".join(current))
            current = []
            current_len = 0
        current.append(para)
        current_len += len(para)

    if current:
        chunks.append("\n\n".join(current))

    return chunks if chunks else [text]
