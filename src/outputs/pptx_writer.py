"""PPTX 输出适配器 - 学术简洁风 PowerPoint 演示文稿."""

from __future__ import annotations

import re
from pathlib import Path

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from src.core.models import Report

# ── Color Palette ──────────────────────────────────────────────────────────
BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_ACCENT = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x99, 0x99, 0x99)

# ── Layout Constants (16:9) ────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
TITLE_BAR_H = Inches(1.0)
FOOTER_H = Inches(0.4)
CONTENT_TOP = TITLE_BAR_H + Inches(0.3)
CONTENT_H = SLIDE_H - CONTENT_TOP - FOOTER_H - Inches(0.2)

# ── Content Splitting ──────────────────────────────────────────────────────
MAX_BULLETS_PER_SLIDE = 10
MAX_CHARS_PER_SLIDE = 1200

FONT_NAME = "Calibri"


def write_pptx(report: Report, output_path: str) -> str:
    """将报告写入 PPTX 文件.

    Args:
        report: 报告对象
        output_path: 输出文件路径

    Returns:
        实际写入的文件路径
    """
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    _add_title_slide(
        prs,
        report.title,
        f"{report.generated_at.strftime('%Y-%m-%d %H:%M')}  |  {', '.join(report.source_files)}",
    )

    # Parse markdown into slides
    _parse_markdown(prs, report.content)

    prs.save(output_path)
    logger.info(f"PPTX report written to: {path}")
    return str(path)


# ── Slide Builders ─────────────────────────────────────────────────────────


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Cover page: full dark-blue background, white centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Full-screen dark blue background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BLUE_DARK

    # Title
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(2.4), CONTENT_W, Inches(1.2))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_NAME

    # Divider line
    line_top = Inches(3.8)
    line_w = Inches(4)
    line_left = (SLIDE_W - line_w) // 2
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        line_left, line_top, line_w, Pt(1),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # Subtitle
    txbox2 = slide.shapes.add_textbox(MARGIN_L, Inches(4.1), CONTENT_W, Inches(0.8))
    tf2 = txbox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(16)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME


def _add_section_header_slide(prs: Presentation, title: str) -> None:
    """Section divider: white background, left dark-blue vertical bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Left vertical bar
    bar = slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(0), Emu(0), Inches(0.35), SLIDE_H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()

    # Section title
    txbox = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), CONTENT_W, Inches(1.2))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = BLUE_DARK
    run.font.name = FONT_NAME


def _add_content_slide(
    prs: Presentation, title: str, items: list[str], page_num: int,
) -> None:
    """Content slide: dark-blue title bar + white body with bulleted items."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title bar background
    bar = slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(0), Emu(0), SLIDE_W, TITLE_BAR_H,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE_DARK
    bar.line.fill.background()

    # Title text
    txbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.2), CONTENT_W, Inches(0.6))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = FONT_NAME

    # Content body
    body_box = slide.shapes.add_textbox(
        MARGIN_L, CONTENT_TOP, CONTENT_W, CONTENT_H,
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.auto_size = None

    for i, item in enumerate(items):
        is_bullet = item.startswith("- ") or item.startswith("* ")
        text = item[2:] if is_bullet else item

        if i == 0:
            para = body_tf.paragraphs[0]
        else:
            para = body_tf.add_paragraph()

        if is_bullet:
            para.level = 0
            para.space_before = Pt(4)
            # Add bullet character
            _add_formatted_text(para, "\u2022  " + text, is_bullet=True)
        else:
            para.space_before = Pt(6)
            _add_formatted_text(para, text, is_bullet=False)

    # Footer: page number
    footer = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.5), SLIDE_H - FOOTER_H, Inches(1.0), FOOTER_H,
    )
    ft = footer.text_frame
    ft.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run = ft.paragraphs[0].add_run()
    run.text = str(page_num)
    run.font.size = Pt(10)
    run.font.color.rgb = GRAY
    run.font.name = FONT_NAME


def _add_formatted_text(para, text: str, is_bullet: bool = False) -> None:
    """Add text with inline **bold** and *italic* formatting via runs."""
    font_size = Pt(14)
    # Split on **bold** and *italic* patterns
    parts = re.split(r"(\*\*.*?\*\*|\*.*?\*)", text)

    for part in parts:
        if not part:
            continue
        run = para.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("*") and part.endswith("*"):
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part

        run.font.size = font_size
        run.font.color.rgb = TEXT_COLOR
        run.font.name = FONT_NAME


# ── Markdown Parser ────────────────────────────────────────────────────────


def _parse_markdown(prs: Presentation, markdown_text: str) -> None:
    """Parse Markdown content and build slides.

    ## heading → section header slide + content slides
    ### heading → content slide(s)
    Lists/paragraphs → fill current content slide
    """
    lines = markdown_text.split("\n")
    current_title = ""
    current_items: list[str] = []
    page_num = 1  # cover is unnumbered

    def flush_content():
        nonlocal page_num
        if not current_title and not current_items:
            return
        if not current_items:
            return
        # Split into pages
        for chunk in _split_items(current_items):
            page_num += 1
            slide_title = current_title
            if page_num > 2 and len(_split_items(current_items)) > 1:
                # Mark continuation only if multi-page
                pass
            _add_content_slide(prs, slide_title, chunk, page_num)

    for line in lines:
        stripped = line.strip()

        # ## heading → section + new content group
        h2_match = re.match(r"^##\s+(.+)$", stripped)
        if h2_match:
            flush_content()
            section_title = h2_match.group(1)
            _add_section_header_slide(prs, section_title)
            page_num += 1
            current_title = section_title
            current_items = []
            continue

        # ### sub-heading → flush + new content group
        h3_match = re.match(r"^###\s+(.+)$", stripped)
        if h3_match:
            flush_content()
            current_title = h3_match.group(1)
            current_items = []
            continue

        # # top-level heading → skip (already on cover)
        if re.match(r"^#\s+", stripped):
            continue

        # Separator → skip
        if re.match(r"^---+$", stripped):
            continue

        # Empty lines → skip
        if not stripped:
            continue

        # Content line
        current_items.append(stripped)

    flush_content()


def _split_items(items: list[str]) -> list[list[str]]:
    """Split a list of items into chunks that fit on one slide."""
    chunks: list[list[str]] = []
    current_chunk: list[str] = []
    current_chars = 0

    for item in items:
        item_len = len(item)
        would_exceed_bullets = len(current_chunk) >= MAX_BULLETS_PER_SLIDE
        would_exceed_chars = current_chars + item_len > MAX_CHARS_PER_SLIDE and current_chunk

        if would_exceed_bullets or would_exceed_chars:
            chunks.append(current_chunk)
            current_chunk = []
            current_chars = 0

        current_chunk.append(item)
        current_chars += item_len

    if current_chunk:
        chunks.append(current_chunk)

    return chunks if chunks else [[]]
